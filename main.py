import os
from docx import Document
from pptx import Presentation
from openpyxl import Workbook
from datetime import datetime

import tkinter as tk
from tkinter import messagebox
import shutil


class MakerInterface:
    def __init__(self, sorter_app):
        self.sorter_app = sorter_app
        self.result_folder = tk.StringVar()
   
        
        
    def sorter_make_form(self):
        for widget in self.sorter_app.winfo_children():
            widget.pack_forget()

        self.sorter_app._add_label(text=" Create/ Search", font=('Helvetica', 50), bg='darkgreen', foreground='white',
                                    width=50)
        
        button_txt = ["Microsoft Office", "Notepad", "Visual-Studio-Code","File Manager","File Organize","Back"]
        for text in button_txt:
            new_button = self.sorter_app._add_button(text, foreground='white', font=("Helvetica", 30), width=20, bg='forest green')
            if text == "Microsoft Office":
                new_button.config(command= self.Microsoft_Office)
            elif text == "Notepad":
                new_button.config(command=self.Notepad_form)
            elif text == "Visual-Studio-Code":
                new_button.config(command=self.visual_studio_code)
            elif text == "File Manager":
                new_button.config(command=self.file_manager_interface)
                
            elif text == "File Organize":
                new_button.config(command=self.file_organize_interface)
                
            elif text =="Back":
                new_button.config(command=self.sorter_app.menu_choice)
            new_button.pack( pady=10, padx=80)
            
            
    def file_organize_interface(self):
   
        for widget in self.sorter_app.winfo_children():
            widget.pack_forget()
        
        self.sorter_app._add_label(text=" Organize Folder Path", font=('Helvetica', 50), bg='darkgreen', foreground='white',
                                    width=50)
        
        self.file_folderpath_entry, self.button = self.sorter_app._add_entry_with_placeholder_and_button("Organize Folder", "Folder", font=('Helvetica', 30),width=45, text='↻', command=self.file_oragnize)
        button_configurations = [
        {"text": "Folder Name    ↑", "width": 20, "command": lambda: self.algrorithm_bubblesort_organizer("asc")},
        {"text": "↓", "width": 2, "command": lambda: self.algrorithm_bubblesort_organizer("desc")},
        
        {"text": "Date Modified  ↑", "width": 17, "command": lambda: self.algrorithm_bubblesort_organizer("asc1")},
        {"text": "↓", "width": 2, "command": lambda: self.algrorithm_bubblesort_organizer("desc1")},
        
        {"text": "Size  ↑", "width": 5, "command": lambda: None},
        {"text": "↓", "width": 2, "command": lambda: None}
        ]
        
        button_frame = tk.Frame(self.sorter_app)
        button_frame.pack()

        for config in button_configurations:
            button = tk.Button(button_frame, text=config["text"], font=("Helvetica", 30), width=config["width"], bg='forest green', command=config["command"])
            button.pack(side=tk.LEFT, pady=10)
        
        self.sorter_app._add_listbox_and_scrollbar(text_variable=self.result_folder, font=('Helvetica', 20), bg='Pink', width=76)

        self.sorter_app._add_button("Back" ,font=("Helvetica", 20), width=72 , command=self.sorter_make_form)
        
    def file_oragnize(self):
        try:
            path_folder_name = self.file_folderpath_entry.get()
            user_home = os.path.expanduser('~')
            search_folder = os.path.join(user_home,path_folder_name) 
        
            files = os.listdir(search_folder)
        
            organized = True
        # Define destination folders
            destinations = {
            '.docx': 'Document-Folder',
            '.py': 'Python-Folder',
            '.pptx': 'Powerpoint-Folder',
            '.txt': 'NotePad-Folder',
            '.xlsx': 'Excel-Folder',
            '.cs': 'C#-Folder',
            '.java': 'Java-Folder',
            '.js': 'JavaScript-Folder',
            '.html': 'HTML-Folder',
            '.css': 'CSS-Folder',
            '.json': 'JSON-Folder',
            '.pdf': 'PDF-Folder', 
            '.mp4': 'Video-Folder-mp4',
            '.mp3': 'Audio-Folder-mp3',
            '.png': 'Image-Folder-png',
            '.jpg': 'Image-Folder-jpg',
            '.jpeg': 'Image-Folder-jpeg',
            '.gif': 'Image-Folder-gif',
            '.zip': 'Zip-Folder',
            '.exe': 'Executable-Folder',
        }
            for file in files:
                filename, extension = os.path.splitext(file)
                destination_folder = destinations.get(extension, extension[1:])
                destination_path = os.path.join(search_folder, destination_folder)

                if not os.path.exists(destination_path):
                    os.makedirs(destination_path)
                    organized=False

                shutil.move(os.path.join(search_folder, file), os.path.join(destination_path, file))
        
            if organized:
                messagebox.showinfo("Success", f"The files in {path_folder_name} are already organized.")
            else:
                messagebox.showinfo("Sucessfully", f"Kindly See the {path_folder_name}")
        except:
            #Location
            sample_usage = r"C:Users\Name\Next\Desktop\Newfolder"
            messagebox.showerror("Try Again", f" Example:\n {sample_usage}")
            
    def file_manager_interface(self):
        for widget in self.sorter_app.winfo_children():
            widget.pack_forget()
    
        self.file_folder_entry, self.button = self.sorter_app._add_entry_with_placeholder_and_button("Search Folder", "Folder", font=('Helvetica', 30),width=45, text='Search', command=self.file_folder)
        
        button_configurations = [
        {"text": "Folder Name    ↑", "width": 20, "command": lambda: self.algrorithm_bubblesort("asc")},
        {"text": "↓", "width": 2, "command": lambda: self.algrorithm_bubblesort("desc")},
        {"text": "Date Modified  ↑", "width": 17, "command": lambda: self.algrorithm_bubblesort("asc1")},
        {"text": "↓", "width": 2, "command": lambda: self.algrorithm_bubblesort("desc1")},
        {"text": "Size  ↑", "width": 5, "command": lambda: None},
        {"text": "↓", "width": 2, "command": lambda: None}
        ]

        button_frame = tk.Frame(self.sorter_app)
        button_frame.pack()
        for config in button_configurations:
            button = tk.Button(button_frame, text=config["text"], font=("Helvetica", 30), width=config["width"], bg='forest green', command=config["command"])
            button.pack(side=tk.LEFT, pady=10)

        self.result_listbox = self.sorter_app._add_listbox_and_scrollbar(text_variable=self.result_folder, font=('Helvetica', 20), bg='Pink', width=76)
        self.sorter_app._add_button("Back" ,font=("Helvetica", 20), width=72 , command=self.sorter_make_form)
        
    def algrorithm_bubblesort_organizer(self, changing):
        path_folder_name = self.file_folderpath_entry.get()
        user_home = os.path.expanduser('~')
        search_folder = os.path.join(user_home,path_folder_name)  
        folders = [folder for folder in os.listdir(search_folder) if os.path.isdir(os.path.join(search_folder, folder))]
        
        if changing == "asc":
            folders = self.operation_bubbleSort(folders)
            
        elif changing == "desc": 
            folders = self.operation_bubbleSortreversed(folders)
        
        elif changing == "asc1":
            folders = self.operation_bubbleSort(folders, key=lambda folder: os.path.getmtime(os.path.join(search_folder, folder)))
        elif changing == "desc1":
            folders = self.operation_bubbleSortreversed(folders, key=lambda folder: os.path.getmtime(os.path.join(search_folder, folder)))

            
        folder_info = []
        for folder_name in folders:
            folder_path = os.path.join(search_folder, folder_name)
            last_modified = os.path.getmtime(folder_path)
            last_modified_formatted = datetime.fromtimestamp(last_modified).strftime("%m/%d/%Y %I:%M%p")
            
            format_folder_style = f'{folder_name:<10}{last_modified_formatted:>80}'
            folder_info.append(format_folder_style)
            
        if folder_info:
            filtered_folders = [f'  "{info}"' for info in folder_info]
            self.result_folder.set("\n".join(filtered_folders))
        else:
            self.result_folder.set("No folders found in the specified directory.")
    #ALGORITHIMS
    def operation_Folder_search(self, s1, s2):
        len_string1 = len(s1)
        len_string2 = len(s2)
        
        dp = []
        row = 0
        while row <= len_string1:
            dp.append([0] * (len_string2 + 1))
            row += 1
        
        i = 0
        while i < len_string1 + 1:
            dp[i][0] = i
            i += 1

        j = 0
        while j < len_string2 + 1:
            dp[0][j] = j
            j += 1

        i = 1
        while i < len_string1 + 1:
            j = 1
            while j < len_string2 + 1:
                cost = 0 if s1[i - 1] == s2[j - 1] else 1
        
                print(cost)
                dp[i][j] = min(dp[i - 1][j] + 1,  # deletion
                               dp[i][j - 1] + 1,  # insertion
                               dp[i - 1][j - 1] + cost)  # substitution
                j += 1
            i += 1
        return dp[len_string1][len_string2]
    
    def operation_bubbleSort(self, array, key= None):
        
        n = len(array)
        while n > 0:
            i = 0
            while i < n - 1:
                if key is not None:
                    if key(array[i].lower()) > key(array[i + 1].lower()):
                        temp = array[i]
                        array[i] = array[i + 1]
                        array[i + 1] = temp 
                
                else:
        
                    if array[i].lower() > array[i + 1].lower():
                        temp = array[i]
                        array[i] = array[i + 1]
                        array[i + 1] = temp 
                i += 1
            n -= 1 
        return array
    
    def operation_bubbleSortreversed(self, array , key= None):

        n = len(array)
        while n > 0:
            i = 0
            while i < n - 1:
                if key is not None:
                    if key(array[i].lower()) < key(array[i + 1].lower()):
                        temp = array[i]
                        array[i] = array[i + 1]
                        array[i + 1] = temp 
                else:
                    if array[i].lower() < array[i + 1].lower():  
                        temp = array[i]
                        array[i] = array[i + 1]
                        array[i + 1] = temp 
                i += 1
            n -= 1 
        return array

    def file_folder(self):
        get_folder = self.file_folder_entry.get()
        if not get_folder:
            messagebox.showerror("Error", "You must enter a folder name.")
            return

        search_folder = r"C:Users\Name\Next\Desktop\Newfolder"
        folders = os.listdir(search_folder)

        closest_match = []

        for folder in folders:
            if get_folder.lower() in folder.lower():
                folder_name, extension = os.path.splitext(folder)
                distance = self.operation_Folder_search(get_folder, folder_name)
                closest_match.append((folder_name, extension, distance))
                
        closest_match.sort(key=lambda x: x[2])
        
        if closest_match:
            formatted_folders = ['"' + folder[0] + folder[1] + '"' for folder in closest_match]
            self.result_folder.set("\n".join(formatted_folders))
        else:
            self.result_folder.set(f"No folders found containing the characters '{get_folder}'.")
            
    def algrorithm_bubblesort(self, changing):
        
        user_home = os.path.expanduser('~')
        
        search_folder = os.path.join(user_home, "OneDrive", "Desktop","File-Sorter")  
        folders = [folder for folder in os.listdir(search_folder) if os.path.isdir(os.path.join(search_folder, folder))]
        print(folders)
        
        #Sorting Algorithm Folder Asceneding And Descending
        if changing == "asc":
            folders = self.operation_bubbleSort(folders)
            
        elif changing == "desc": 
            folders = self.operation_bubbleSortreversed(folders)
            
        #Date modified Sort by New to old date
            # Sorting by last modified date
        elif changing == "asc1":
            folders = self.operation_bubbleSort(folders, key=lambda folder: os.path.getmtime(os.path.join(search_folder, folder)))
        elif changing == "desc1":
            folders = self.operation_bubbleSortreversed(folders, key=lambda folder: os.path.getmtime(os.path.join(search_folder, folder)))

        folder_info = []
        for folder_name in folders:
            folder_path = os.path.join(search_folder, folder_name)
            last_modified = os.path.getmtime(folder_path)
            last_modified_formatted = datetime.fromtimestamp(last_modified).strftime("%m/%d/%Y %I:%M%p")
            
            format_folder_style = f'{folder_name:<10}{last_modified_formatted:>80}'
            folder_info.append(format_folder_style)
            
        if folder_info:
            filtered_folders = [f'  "{info}"' for info in folder_info]
            self.result_folder.set("\n".join(filtered_folders))
        else:
            self.result_folder.set("No folders found in the specified directory.") 
        
    def Microsoft_Office(self):
        for widget in self.sorter_app.winfo_children():
            widget.pack_forget()

        self.sorter_app._add_label(text="Microsoft Office", font=('Helvetica', 50), bg='darkgreen', foreground='white',width=50)
        button_txt = ["Powerpoint", "Document", "Excel", "Back"]
        for text in button_txt:
            new_button = self.sorter_app._add_button(text, foreground='white', font=("Helvetica", 30), width=30, bg='forest green')
            
            if text == "Powerpoint":
                new_button.config(command=lambda:self.Microsoft_Office_File('Powerpoint'))                
            elif text == "Document":
                new_button.config(command=lambda:self.Microsoft_Office_File('Document'))
            elif text == "Excel":
                new_button.config(command=lambda:self.Microsoft_Office_File('Excel'))
            elif text =="Back":
                new_button.config(command=self.sorter_make_form)
            new_button.pack(side=tk.RIGHT, pady=10, padx=80)
            
    def Microsoft_Office_File(self, file_type):
        for widget in self.sorter_app.winfo_children():
            widget.pack_forget()
        self.sorter_app._add_label(text=f"Microsoft Office: {file_type}", font=('Helvetica', 50), bg='darkgreen', foreground='white',width=50)
        
        if file_type == 'Powerpoint':
            self.file_folder_entry = self.sorter_app._add_entry_with_placeholder("Powerpoint-Folder", "Powerpoint", font=('Helvitica', 50))
            self.file_name_entry = self.sorter_app._add_entry_with_placeholder("Powerpoint File Name", "Powerpoint", font=('Helvitica', 50))
            
        elif file_type == 'Document':
            self.file_folder_entry = self.sorter_app._add_entry_with_placeholder("Document-Folder", "Document", font=('Helvitica', 50))
            self.file_name_entry = self.sorter_app._add_entry_with_placeholder("Document File Name", "Document", font=('Helvitica', 50))
        
        elif file_type == "Excel":
            self.file_folder_entry = self.sorter_app._add_entry_with_placeholder("Excel-Folder", "Excel", font=('Helvitica', 50))
            self.file_name_entry = self.sorter_app._add_entry_with_placeholder("Excel File Name", "Excel", font=('Helvitica', 50))
            
        else:
            None
        button_txt = [f"Create {file_type} File", "Back"]
        for text in button_txt:
            new_button = self.sorter_app._add_button(text, foreground='white', font=("Helvetica", 30), width=25, bg='forest green')
            if text == f"Create {file_type} File":
                new_button.config(command=lambda:self.create_microsoft_file(file_type,self.file_folder_entry.get(), self.file_name_entry.get()))
            elif text =="Back":
                new_button.config(command=self.sorter_make_form)
            new_button.pack(side=tk.RIGHT, pady=15, padx=80)
            
    def create_microsoft_file(self, file_type, file_folder, file_name):
        try:
            folder_name = file_folder
            micro_file_name = file_name
            if file_type == 'Powerpoint':
                validation_file_type = ["Document-Folder", "Excel-Folder"]
                if any(value == "" for value in [micro_file_name, folder_name]) :
                    messagebox.showerror("Try Again", 'You cannot create \nEnter Filename to Create')           
                
                elif folder_name in validation_file_type:
                    messagebox.showerror("Error","Cannot Save In This Folder") 

                elif os.path.exists(os.path.join(folder_name, micro_file_name + ".pptx")):
                    result = messagebox.askokcancel("File existing", f"Do you want to open in Microsoft Office\nOpen {micro_file_name}.pptx in Microsoft Office?")
                    if result:
                        os.system(f'start {folder_name}\\{micro_file_name}.pptx')
                    else:
                        return
                else:
                    prs = Presentation()
                    slide_layout = prs.slide_layouts[0]
                    slide = prs.slides.add_slide(slide_layout)
                    title = slide.shapes.title 
                    title.text = 'My Presentation' 
                    result = messagebox.askokcancel("Do you want to create", f"Powerpoint {micro_file_name}")
                    if result:
                        prs.save(micro_file_name + '.pptx')
                        os.makedirs(folder_name, exist_ok=True)
                        os.rename(micro_file_name + '.pptx', os.path.join(folder_name, micro_file_name + '.pptx'))
                        os.system(f'start {folder_name}\\{micro_file_name}.pptx')
                    else:
                        return
                      
            elif file_type == "Document":
                validation_file_type = ["Powerpoint-Folder", "Excel-Folder"]
                if any(value == "" for value in [micro_file_name, folder_name]) :
                    messagebox.showerror("Try Again", 'You cannot create \nEnter Filename to Create') 
                    
                elif folder_name in validation_file_type:
                    messagebox.showerror("Error","Cannot Save In This Folder") 
              
                elif os.path.exists(os.path.join(folder_name, micro_file_name + ".pptx")):
                    result = messagebox.askokcancel("File existing", f"Do you want to open in Microsoft Office\nOpen {micro_file_name}.docx in Microsoft Office?")
                    if result:
                        os.system(f'start {folder_name}\\{micro_file_name}.docx')
                    else:
                        return
                else:
                    doc = Document()
                    doc.add_heading('My Document', level=1) 
                    doc.add_paragraph('This is First Pharagraph')
                    result = messagebox.askokcancel("Do you want to create", f"Document {micro_file_name}")
                    if result:
                        doc.save(micro_file_name + ".docx")
                        os.makedirs(folder_name, exist_ok=True)
                        os.rename(micro_file_name + '.docx', os.path.join(folder_name, micro_file_name + '.docx'))
                        os.system(f'start {folder_name}\\{micro_file_name}.docx')
                    else:
                        return
                    
            elif file_type == "Excel":
                validation_file_type = ["Document-Folder", "Presentation-Folder"]
                if any(value == "" for value in [micro_file_name, folder_name]) :
                    messagebox.showerror("Try Again", 'You cannot create \nEnter Filename to Create')
                    
                elif folder_name in validation_file_type:
                    messagebox.showerror("Error","Cannot Save In This Folder") 

                elif os.path.exists(os.path.join(folder_name, micro_file_name + ".pptx")):
                    result = messagebox.askokcancel("File existing", f"Do you want to open in Microsoft Office\nOpen {micro_file_name}.docx in Microsoft Office?")
                    if result:
                        os.system(f'start {folder_name}\\{micro_file_name}.docx')
                    else:
                        return
                else:
                    wb = Workbook()
                    ws = wb.active
                    ws['A1']="Hello"
                    ws['B1']="Finals"
                    result = messagebox.askokcancel("Do you want to create", f"Excel {file_name}")
                    if result:
                        wb.save(file_name + ".xlsx") 
                        os.makedirs(folder_name, exist_ok=True)
                        os.rename(file_name + '.xlsx', os.path.join(folder_name, file_name + '.xlsx'))
                        os.system(f'start {folder_name}\\{file_name}.xlsx')
                    else:
                        return
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred: {str(e)}")
          
    def Notepad_form(self):
        for widget in self.sorter_app.winfo_children():
            widget.pack_forget()
        self.sorter_app._add_label(text="Notepad", font=('Helvetica', 50), bg='darkgreen', foreground='white',width=50)            
        
        self.file_folder_entry = self.sorter_app._add_entry_with_placeholder("Notepad-Folder", "Notepad", font=('Helvitica', 50))
        self.file_name_entry = self.sorter_app._add_entry_with_placeholder("Notepad File Name", "Notepad", font=('Helvitica', 50))
        
        self.sorter_app._add_button(text="Create Notepad  File",  foreground='white', command=self.create_notepad, font=("Helvetica", 30), width=20, bg='forest green')
        self.sorter_app._add_button(text="Back",  foreground='white', command=self.sorter_make_form, font=("Helvetica", 30), width=20, bg='forest green')        

    def create_notepad(self):
        folder_name = self.file_folder_entry.get()
        file_name = self.file_name_entry.get()
        
        if any(value == "" for value in [file_name, folder_name]):
            messagebox.showerror("Try Again", 'You cannot create \nEnter Filename to Create')
        elif os.path.exists(os.path.join(folder_name, file_name + ".txt")):
            result = messagebox.askokcancel("File existing", f"Do you want to open in Text Pad\nOpen {file_name}.txt in Note ?")
            if result:
                os.system(f'start {folder_name}\\{file_name}.txt')
            else:
                return
        else:
            os.makedirs(folder_name, exist_ok=True)
            with open(os.path.join(folder_name, file_name + '.txt'),'w') as f:
                f.write("Your Content Here")
            result = messagebox.askokcancel("Do you want to create", f"Notepad {file_name}")
            if result:
                os.system(f'start {folder_name}\\{file_name}.txt')
            else:
                return
        
    def visual_studio_code(self):
        for widget in self.sorter_app.winfo_children():
            widget.pack_forget()
        self.sorter_app._add_label(text="Programming Language", font=('Helvetica', 50), bg='darkgreen', foreground='white',width=50)
        
        button_txt = ["Python","JavaScript","Java", "C#", "Back"]
        for text in button_txt:
            new_button = self.sorter_app._add_button(text, foreground='white', font=("Helvetica", 30), width=20, bg='forest green')
            if text == "Python":
                new_button.config(command=lambda:self.visual_studio_form('Python', ))
            elif text == "JavaScript":  
                new_button.config(command=lambda:self.visual_studio_form('JavaScript'))
            elif text == "Java":
                new_button.config(command=lambda:self.visual_studio_form('Java'))
            elif text == "C#":
                new_button.config(command=lambda:self.visual_studio_form('C#'))
            elif text == "Back":
                new_button.config(command=self.sorter_make_form)
            new_button.pack(side=tk.RIGHT, pady=10, padx=80) 
            
    def visual_studio_form(self, prog_language):
        for widget in self.sorter_app.winfo_children():
            widget.pack_forget()
        self.sorter_app._add_label(text=f"Programming Language: {prog_language}", font=('Helvetica', 50), bg='darkgreen', foreground='white',width=50)
        
        if prog_language == 'Python':
            self.file_folder_entry = self.sorter_app._add_entry_with_placeholder("Python-Folder", "Python", font=('Helvitica', 50))
            self.file_name_entry = self.sorter_app._add_entry_with_placeholder("Python File Name", "Python", font=('Helvitica', 50))
            
        elif prog_language == 'JavaScript':
            self.file_folder_entry = self.sorter_app._add_entry_with_placeholder("JavaScript-Folder", "JavaScript", font=('Helvitica', 50))
            self.file_name_entry = self.sorter_app._add_entry_with_placeholder("JavaScript File Name", "JavaScript", font=('Helvitica', 50))
        
        elif prog_language == "Java":
            self.file_folder_entry = self.sorter_app._add_entry_with_placeholder("Java-Folder", "Java", font=('Helvitica', 50))
            self.file_name_entry = self.sorter_app._add_entry_with_placeholder("Java File Name", "Java", font=('Helvitica', 50))
            
        elif prog_language == "C#":
            self.file_folder_entry = self.sorter_app._add_entry_with_placeholder("C#-Folder", "C#", font=('Helvitica', 50))
            self.file_name_entry = self.sorter_app._add_entry_with_placeholder("C# File Name", "C#", font=('Helvitica', 50))
        else:
            None
        button_txt = [f"Create {prog_language} File", "Back"]
        for text in button_txt:
            new_button = self.sorter_app._add_button(text, foreground='white', font=("Helvetica", 30), width=25, bg='forest green')
            if text == f"Create {prog_language} File":
                new_button.config(command=lambda:self.create_prog_language_code(prog_language,self.file_folder_entry.get(), self.file_name_entry.get()))
            elif text =="Back":
                new_button.config(command=self.sorter_make_form)
            new_button.pack(side=tk.RIGHT, pady=15, padx=80)
    
    def create_prog_language_code(self, prog_language, file_folder, file_name):
        try:
            folder_name = file_folder
            prog_file_name = file_name
            if prog_language == 'Python':
                if any(value == "" for value in [prog_file_name, folder_name]):
                    messagebox.showerror("Try Again", 'You cannot create \nEnter Filename to Create')
                    
                elif folder_name in ["JavaScript-Folder", "Java-Folder", "C#-Folder"]:
                    messagebox.showwarning("Warning","This feature is not available\nPlease use the appropriate File Folder language")
                    
                elif os.path.exists(os.path.join(folder_name, prog_file_name +  ".py")):
                    result = messagebox.askokcancel("Do you want to open in Visual Studio Code", f"Open {prog_file_name}.py in Visual Studio Code?")
                    if result:
                        os.system(f'start {folder_name}\\{file_name}.py')
                    else:
                        return
                else:
                    os.makedirs(folder_name, exist_ok=True)
                    with open(os.path.join(folder_name, prog_file_name + '.py'),'w') as f:
                        f.write("print(Your Content Here)")
                    result = messagebox.askokcancel("Do you want to create in Visual Studio Code", f"Create {prog_file_name}.py in Visual Studio Code?")
                    if result:
                        os.system(f'start {folder_name}\\{file_name}.py')
                    else:
                        return
            
            elif prog_language == 'JavaScript':
                if any(value == "" for value in [prog_file_name, folder_name]):
                    messagebox.showerror("Try Again", 'You cannot create \nEnter Filename to Create')
                    
                elif folder_name in ["Python-Folder", "Java-Folder", "C#-Folder"]:
                    messagebox.showwarning("Warning","This feature is not available\nPlease use the appropriate File Folder language")
                    
                elif os.path.exists(os.path.join(folder_name, prog_file_name + ".js")):
                    result = messagebox.askokcancel("Do you want to open in Visual Studio Code", f"Open {prog_file_name}.js in Visual Studio Code?")
                    if result:
                        os.system(f'start {folder_name}\\{file_name}.js')
                    else:
                        return
                else:
                    os.makedirs(folder_name, exist_ok=True)
                    with open(os.path.join(folder_name, prog_file_name + '.js'),'w') as f:
                        f.write("console.log(\"Hello World!\");")
                    result = messagebox.askokcancel("Do you want to create in Visual Studio Code", f"Create {prog_file_name}.js in Visual Studio Code?")
                    if result:
                        os.system(f'start {folder_name}\\{file_name}.js')
                    else:
                        return
                    
            elif prog_language == 'Java':
                if any(value == "" for value in [prog_file_name, folder_name]):
                    messagebox.showerror("Try Again", 'You cannot create \nEnter Filename to Create')
                elif folder_name in ["Python-Folder", "JavaScript-Folder", "C#-Folder"]:
                    messagebox.showwarning("Warning","This feature is not available\nPlease use the appropriate File Folder language")
                    
                elif os.path.exists(os.path.join(folder_name, prog_file_name + ".java")):
                    result = messagebox.askokcancel("Do you want to open in Visual Studio Code", f"Open {prog_file_name}.java in Visual Studio Code?")
                    if result:
                        os.system(f'start {folder_name}\\{file_name}.java')
                    else:
                        return 
                else:
                    os.makedirs(folder_name, exist_ok=True)
                    with open(os.path.join(folder_name, prog_file_name + '.java'),'w') as f:
                        f.write('public class ' + prog_file_name + ' {\n')
                        f.write('    public static void main(String[] args) {\n')
                        f.write('        System.out.println("Hello, world!");\n')
                        f.write('    }\n')
                        f.write('}\n')
                    result = messagebox.askokcancel("Do you want to create in Visual Studio Code", f"Create {prog_file_name}.java in Visual Studio Code?")
                    if result:
                        os.system(f'start {folder_name}\\{file_name}.java')
                    else:
                        return
                        
            elif prog_language == 'C#':
                if any(value == "" for value in [prog_file_name, folder_name]):
                    messagebox.showerror("Try Again", 'You cannot create \nEnter Filename to Create')
                    
                elif folder_name in ["Python-Folder", "JavaScript-Folder", "Java-Folder"]:
                    messagebox.showwarning("Warning","This feature is not available\nPlease use the appropriate File Folder language")
                    
                elif os.path.exists(os.path.join(folder_name, prog_file_name + ".cs")):
                    result = messagebox.askokcancel("Do you want to open in Visual Studio Code", f"Open {prog_file_name}.cs in Visual Studio Code?")
                    if result:
                        os.system(f'start {folder_name}\\{file_name}.cs')
                    else:
                        return
                else:
                    os.makedirs(folder_name, exist_ok=True) 
                    with open(os.path.join(folder_name, prog_file_name + '.cs'),'w') as f:
                        f.write('using System;\n\n')
                        f.write('class ' + prog_file_name + ' {\n')
                        f.write('    static void Main() {\n')
                        f.write('        Console.WriteLine("Hello, world!");\n')
                        f.write('    }\n')
                        f.write('}\n')
                    result = messagebox.askokcancel("Do you want to create in Visual Studio Code", f"Create {prog_file_name}.cs in Visual Studio Code?")
                    if result:
                        os.system(f'start {folder_name}\\{file_name}.cs')
                    else:
                        return
                    
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred: {str(e)}")
            
class ShowMessage: 
    def __init__(self, message) -> None:
        self.__message = message

    def __call__(self):
        messagebox.showinfo(title="Acknowledgement", message=self.__message)

class Make_File_App_Builder:
    def __init__(self):
        self.sorter_app = Maker_File_App()

    def build_menu_choice(self):
        self.sorter_app.menu_choice()

    def build_display(self):
        self.sorter_app.display()

class Maker_File_App_Director:
    def __init__(self, builder):
        self.builder = builder

    def construct(self):
        self.builder.build_menu_choice()
        self.builder.build_display()

class Maker_File_App(tk.Tk):
    def __init__(self):
        super().__init__()
        self.makefile_user_interface = MakerInterface(self)
        self.__configure_window()

    def __configure_window(self):
        self.frame()
        self.title("File Creation App")
        self.geometry('1375x750')
        self.resizable(True, False)
        self.configure(bg="green")

    def frame(self):
        pass

    def _add_entry_with_placeholder(self, placeholder, default_text, font=None, width=None):
        entry = PlaceholderEntry(self, placeholder, default_text, font=font, width=width)
        entry.pack(padx=100, pady=30)
        return entry
    
    def _add_entry_with_placeholder_and_button(self, placeholder, default_text, font=None, width=None,text=None,  command=None, state=None):
        frame = tk.Frame(self)
        frame.pack(padx=50, pady=30)
        
        entry = PlaceholderEntry(frame, placeholder, default_text, font=font, width=width)
        entry.pack(side=tk.LEFT)
        
        button = tk.Button(frame, text=text, command=command,state=state , font=font)
        button.pack(side=tk.RIGHT, padx=10, pady=10)
        
        return entry, button
    
    def _add_button(self, text=None, foreground=None, command=None, state=None, font=None,image=None, width=None, bg=None, padx=None,
                    pady=None):
        frame = tk.Frame(self, bg="forest green")
        frame.pack()
        button = tk.Button(frame, font=font, foreground=foreground, text=text, image=image, width=width, command=command,padx=padx, bg=bg, pady=pady, state=state)
        button.pack()
        button.image=image
        return button

    def _add_label(self, textvariable=None, text=None, font=None, bg=None, width=None, foreground=None, padx=None):
        label = tk.Label(self, textvariable=textvariable, text=text, bg=bg, width=width, padx=padx, font=font,foreground=foreground)
        label.pack()
        return label

    def _add_listbox_and_scrollbar(self, text_variable=None, font=None, bg=None, **kwargs):
        listbox_frame = tk.Frame(self)
        listbox_frame.pack()
      
        scrollbar = tk.Scrollbar(listbox_frame, orient="vertical")
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        listbox = tk.Listbox(listbox_frame, yscrollcommand=scrollbar.set, font=font, bg=bg, **kwargs)
        listbox.pack(side="left", fill="both", expand=True)

        scrollbar.config(command=listbox.yview)

        if text_variable:
            listbox.config(listvariable=text_variable)
        return listbox
      
    def menu_choice(self):
        for widget in self.makefile_user_interface.sorter_app.winfo_children():
            widget.forget()
        self._add_label(text="File-Create\n Extension\n Organizer", font=('Helvetica', 40), bg='green', width=20).pack(side=tk.LEFT, padx=40)

        button_txt = ["File-Create Organize", "Close"]
        for text in button_txt:
            new_button = self._add_button(text, foreground='white', font=("Helvetica", 30), width=40, bg='forest green')
            if text == "File-Create Organize":
                new_button.config(command=self.makefile_user_interface.sorter_make_form)
            elif text == "Close":
                new_button.config(command=self.destroy)

            new_button.pack(side=tk.LEFT, pady=100, padx=80)

        self._add_button("About This", command=ShowMessage(
            'Ready To Final Application\nBSCS 2-A\nFaborada Nathaniel\n John Paul Bodino'), foreground='white',
                         font=("Helvetica", 30), width=40, bg='forest green').pack(side=tk.LEFT, pady=60, padx=80)

    def display(self):
        self.mainloop()
        
class PlaceholderEntry(tk.Entry):
    def __init__(self, master=None, placeholder="", default_text="", *args, **kwargs):
        super().__init__(master, *args, **kwargs)
        self.placeholder = placeholder
        self.default_text = default_text
        self.placeholder_shown = True
        self.bind("<FocusIn>", self.focus_in)
        self.bind("<FocusOut>", self.focus_out)
        self.show_placeholder()

    def focus_in(self, _):
        if self.placeholder_shown:
            self.delete("0", "end")
            self.placeholder_shown = False

    def focus_out(self, _):
        if not self.get():
            self.show_placeholder()

    def show_placeholder(self):
        self.insert("0", self.placeholder)
        self.placeholder_shown = True
        self.configure(fg="grey")

if __name__ == "__main__":

    builder = Make_File_App_Builder()
    director = Maker_File_App_Director(builder)
    director.construct()
